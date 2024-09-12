import io
import os
from typing import IO, Any
import csv

import chardet
import docx
import openpyxl
import pptx
from pypdf import PdfReader
from pypdf.errors import PdfStreamError

TEXT_SECTION_SEPARATOR = "\n\n"

VALID_FILE_EXTENSIONS = [
    ".txt", ".md", ".pdf", ".docx", ".pptx", ".xlsx", ".csv"
]

def get_file_ext(file_path_or_name: str) -> str:
    _, extension = os.path.splitext(file_path_or_name)
    return extension.lower()

def check_file_ext_is_valid(ext: str) -> bool:
    return ext in VALID_FILE_EXTENSIONS

def detect_encoding(file: IO[bytes]) -> str:
    raw_data = file.read(50000)
    encoding = chardet.detect(raw_data)["encoding"] or "utf-8"
    file.seek(0)
    return encoding

def read_text_file(file: IO, encoding: str = "utf-8", errors: str = "replace") -> str:
    file_content_raw = ""
    for line in file:
        try:
            line = line.decode(encoding) if isinstance(line, bytes) else line
        except UnicodeDecodeError:
            line = line.decode(encoding, errors=errors) if isinstance(line, bytes) else line
        file_content_raw += line
    return file_content_raw

def pdf_to_text(file: IO[Any]) -> str:
    try:
        pdf_reader = PdfReader(file)
        return TEXT_SECTION_SEPARATOR.join(
            page.extract_text() for page in pdf_reader.pages
        )
    except PdfStreamError:
        print("PDF file is not a valid PDF")
    except Exception:
        print("Failed to read PDF")
    return ""

def docx_to_text(file: IO[Any]) -> str:
    doc = docx.Document(file)
    full_text = [para.text for para in doc.paragraphs]
    return TEXT_SECTION_SEPARATOR.join(full_text)

def pptx_to_text(file: IO[Any]) -> str:
    presentation = pptx.Presentation(file)
    text_content = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        extracted_text = f"\nSlide {slide_number}:\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                extracted_text += shape.text + "\n"
        text_content.append(extracted_text)
    return TEXT_SECTION_SEPARATOR.join(text_content)

def xlsx_to_text(file: IO[Any]) -> str:
    workbook = openpyxl.load_workbook(file)
    text_content = []
    for sheet in workbook.worksheets:
        sheet_string = "\n".join(
            ",".join(map(str, row))
            for row in sheet.iter_rows(min_row=1, values_only=True)
        )
        text_content.append(sheet_string)
    return TEXT_SECTION_SEPARATOR.join(text_content)

def csv_to_text(file: IO[Any]) -> str:
    csv_content = file.read().decode('utf-8')
    csv_reader = csv.reader(io.StringIO(csv_content))
    return "\n".join([", ".join(row) for row in csv_reader])

def extract_file_text(file_name: str, file: IO[Any]) -> str:
    extension = get_file_ext(file_name)
    if not check_file_ext_is_valid(extension):
        raise ValueError(f"Unsupported file extension: {extension}")

    if extension in ['.txt', '.md']:
        return read_text_file(file)
    elif extension == '.pdf':
        return pdf_to_text(file)
    elif extension == '.docx':
        return docx_to_text(file)
    elif extension == '.pptx':
        return pptx_to_text(file)
    elif extension in ['.xlsx', '.csv']:
        return csv_to_text(file)
    else:
        raise ValueError(f"Unsupported file extension: {extension}")